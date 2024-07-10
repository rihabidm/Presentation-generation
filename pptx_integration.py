#Importing the necessary libraries

from pptx import Presentation
from pptx.util import Inches
import openai

def generate_presentation_content(content):
    response = openai.Completion.create(
        engine="davinci-codex",
        prompt=f"Generate a detailed presentation based on the following description: {content}",
        max_tokens=500
    )
    return response.choices[0].text.strip()

def create_presentation(content):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] 

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Generated Presentation"
    body.text = content

    return prs

def save_presentation(prs, filename):
    prs.save(filename)